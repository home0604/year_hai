# -*- coding: utf-8 -*-
"""fm-meeting-deck 를 편집 가능한 .pptx 로 생성. python-pptx 네이티브 요소만 사용."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

C = dict(
    bg="FFFFFF", sub="F5F7F8", ink="181C21", ink2="3E464F", muted="78818B",
    line="D5DADF", line2="E6EAED",
    dfm="A85933", dfmbg="F5E7DF", cfm="256B7B", cfmbg="DEECEF",
    ar="5C5471", arbg="E8E5EE", ok="2F6B45", okbg="DFEDE4", white="FFFFFF",
)
SANS = "Malgun Gothic"
MONO = "Consolas"
def clr(h): return RGBColor.from_string(h)
EMU = 914400
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = clr(C["bg"]); bg.line.fill.background()
    bg.shadow.inherit = False
    return s

def _font(run, size, color, bold, font, italic=False):
    run.font.size = Pt(size); run.font.name = font; run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = clr(color)
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', font)

def txt(s, l, t, w, h, lines, size=14, color=None, bold=False, font=SANS,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=4, line_sp=1.05):
    color = color or C["ink2"]
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(lines, str):
        lines = [[(lines, size, color, bold, False, font)]]
    for i, para in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0); p.line_spacing = line_sp
        if isinstance(para, str):
            para = [(para, size, color, bold, False, font)]
        for spec in para:
            tx, sz, cl, bd, it, fn = (list(spec) + [size, color, bold, False, font])[:6]
            r = p.add_run(); r.text = tx
            _font(r, sz, cl or color, bd, fn, it)
    return tb

def rect(s, l, t, w, h, fill=None, line=None, lw=1.0, dash=False):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill: sp.fill.solid(); sp.fill.fore_color.rgb = clr(fill)
    else: sp.fill.background()
    if line:
        sp.line.color.rgb = clr(line); sp.line.width = Pt(lw)
        if dash:
            ln = sp.line._get_or_add_ln(); ln.append(ln.makeelement(qn('a:prstDash'), {'val':'dash'}))
    else:
        sp.line.fill.background()
    return sp

def oval(s, l, t, w, h, fill=None, line=None, lw=1.0):
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill: sp.fill.solid(); sp.fill.fore_color.rgb = clr(fill)
    else: sp.fill.background()
    if line: sp.line.color.rgb = clr(line); sp.line.width = Pt(lw)
    else: sp.line.fill.background()
    return sp

def arrow(s, l, t, w=0.55, size=22):
    txt(s, l, t, w, 0.4, [[("→", size, C["muted"], False, False, SANS)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def kicker_title(s, kicker, title, lede=None, num=None):
    txt(s, 0.62, 0.40, 11, 0.3, [[(kicker.upper(), 11, C["muted"], False, False, MONO)]])
    txt(s, 0.62, 0.70, 11.4, 0.9, [[(title, 24, C["ink"], True, False, SANS)]], line_sp=1.05)
    if num:
        txt(s, 11.9, 0.40, 1.1, 0.3, [[(num, 11, C["muted"], False, False, MONO)]], align=PP_ALIGN.RIGHT)
    y = 1.42
    if lede:
        txt(s, 0.62, 1.44, 11.9, 0.9, [[(lede, 12.5, C["ink2"], False, False, SANS)]], line_sp=1.12)
        y = 2.12
    rect(s, 0.62, y, 12.1, 0.012, fill=C["line2"])
    return y + 0.18

def card(s, l, t, w, h, accent=None, fill=None):
    rect(s, l, t, w, h, fill=fill or C["sub"], line=C["line2"], lw=0.75)
    if accent: rect(s, l, t, w, 0.05, fill=accent)

def _no_style(tbl):
    tblPr = tbl._tbl.tblPr
    for e in tblPr.findall(qn('a:tableStyleId')): tblPr.remove(e)
    st = tblPr.makeelement(qn('a:tableStyleId'), {})
    st.text = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'   # No Style, No Grid
    tblPr.append(st)

def _cell_border(cell, color, w=9525, sides=('bottom',)):
    tcPr = cell._tc.get_or_add_tcPr()
    tag = {'left':'a:lnL','right':'a:lnR','top':'a:lnT','bottom':'a:lnB'}
    order = ['a:lnL','a:lnR','a:lnT','a:lnB']
    for side in sides:
        t_ = tag[side]
        for e in tcPr.findall(qn(t_)): tcPr.remove(e)
        ln = tcPr.makeelement(qn(t_), {'w':str(w),'cap':'flat'})
        fill = ln.makeelement(qn('a:solidFill'), {})
        c_ = fill.makeelement(qn('a:srgbClr'), {'val':color}); fill.append(c_); ln.append(fill)
        # insert in schema order (lnL,lnR,lnT,lnB before fill/other)
        tcPr.append(ln)

def table(s, l, t, w, headers, rows, colw=None, fs=11.5, header_fs=10,
          row_h=0.36, hi=None, win=None):
    nrow, ncol = len(rows)+1, len(headers)
    gt = s.shapes.add_table(nrow, ncol, Inches(l), Inches(t), Inches(w), Inches(row_h*nrow)).table
    gt.first_row = False; gt.horz_banding = False
    _no_style(gt)
    if colw:
        tot = sum(colw)
        for c, cwd in enumerate(colw):
            gt.columns[c].width = Emu(int(w*EMU*cwd/tot))
    def style(cell, align, top=False):
        cell.fill.solid(); cell.fill.fore_color.rgb = clr(C["bg"])
        cell.margin_left = cell.margin_right = Pt(6)
        cell.margin_top = cell.margin_bottom = Pt(2)
        cell.vertical_anchor = MSO_ANCHOR.BOTTOM if top else MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]; p.alignment = align; return p
    for c, h in enumerate(headers):
        cell = gt.cell(0,c)
        p = style(cell, PP_ALIGN.LEFT if c==0 else PP_ALIGN.RIGHT, top=True)
        r = p.add_run(); r.text = h; _font(r, header_fs, C["muted"], True, MONO)
        _cell_border(cell, C["line"], w=12700, sides=('bottom',))
    for ri, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = gt.cell(ri+1,c)
            p = style(cell, PP_ALIGN.LEFT if c==0 else PP_ALIGN.RIGHT)
            col, bold = C["ink2"], False
            fn = SANS if c==0 else MONO
            if win and (ri,c) in win: col, bold = C["ok"], True
            elif hi and (ri,c) in hi: col, bold = C["ink"], True
            r = p.add_run(); r.text = str(val); _font(r, fs, col, bold, fn)
            _cell_border(cell, C["line2"], w=9525, sides=('bottom',))
    return gt

def codebox(s, l, t, w, h, lines, fs=10.5, title=None):
    if title:
        txt(s, l, t, w, 0.3, [[(title,12.5,C["ink"],True,False,SANS)]]); t += 0.34; h -= 0.34
    rect(s, l, t, w, h, fill=C["sub"], line=C["line2"], lw=0.75)
    tb = s.shapes.add_textbox(Inches(l+0.14), Inches(t+0.1), Inches(w-0.28), Inches(h-0.2))
    tf = tb.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after = Pt(0); p.space_before = Pt(0); p.line_spacing = 1.12
        cmt = ln.strip().startswith("#")
        r = p.add_run(); r.text = ln if ln else " "
        _font(r, fs, C["muted"] if cmt else C["ink"], False, MONO)

def box(s, l, t, w, h, lines, accent="ok", fs=12.5):
    bg = {"ok":C["okbg"], "dfm":C["dfmbg"]}[accent]; ac = {"ok":C["ok"], "dfm":C["dfm"]}[accent]
    rect(s, l, t, w, h, fill=bg); rect(s, l, t, 0.05, h, fill=ac)
    txt(s, l+0.22, t+0.14, w-0.4, h-0.24, lines, size=fs, color=C["ink2"], line_sp=1.18, sp_after=5)

def patch(s, l, t, cell, states, outline=None):
    fills = {'m':C["dfm"], 'e':C["sub"], 'x':C["line"]}
    for (cx,cy),st in zip([(0,0),(1,0),(0,1),(1,1)], states):
        r = rect(s, l+cx*cell, t+cy*cell, cell, cell, fill=fills[st],
                 line=(C["muted"] if st=='x' else (C["line"] if st=='e' else None)),
                 lw=0.75, dash=(st=='x'))
    if outline:
        rect(s, l-0.006, t-0.006, cell*2+0.012, cell*2+0.012, line=outline, lw=1.5)

# ======================================================================
# 1. Problem
# ======================================================================
s = slide()
y = kicker_title(s, "Problem", "안테나 마스크 역설계",
    "목표 스펙트럼이 주어지면 그것을 만들어내는 10×10 이진 금속 마스크를 생성한다. 기존 접근은 autoregressive(AR) 모델이다. "
    "Flow model 로 같은 문제를 풀고 AR 과 비교하는 것이 이 발표의 범위다.", num="1 / 10")
fw = 3.9
for i,(a,b,c) in enumerate([("INPUT","y ∈ R^201","목표 S11 스펙트럼 (201 주파수점)"),
    ("GENERATE","b ∈ {0,1}^100","10×10 이진 마스크 — 생성모델이 만든다"),
    ("VERIFY","ŷ = f_fwd(b)","동결된 forward surrogate CNN 으로 재평가")]):
    x = 0.62 + i*(fw+0.12); card(s, x, y, fw, 0.95)
    txt(s, x+0.18, y+0.12, fw-0.36, 0.25, [[(a,10,C["muted"],False,False,MONO)]])
    txt(s, x+0.18, y+0.34, fw-0.36, 0.3, [[(b,13,C["ink"],False,False,MONO)]])
    txt(s, x+0.18, y+0.62, fw-0.36, 0.4, [[(c,11.5,C["ink2"],False,False,SANS)]], line_sp=1.05)
y2 = y+1.15; cw=6.0
card(s, 0.62, y2, cw, 1.95)
txt(s, 0.82, y2+0.14, cw-0.4, 0.3, [[("문제의 성격",13,C["ink"],True,False,SANS)]])
txt(s, 0.82, y2+0.5, cw-0.4, 1.35, [
    "• 정방향(마스크→스펙트럼)은 함수지만 역방향은 one-to-many — 같은 스펙트럼을 내는 마스크가 여러 개.",
    "• 목표는 점추정이 아니라 조건부 분포 p(b|y) 의 샘플링. 생성모델이 필요한 이유.",
    "• 데이터: train 167,430 / test 20,929 쌍."], size=12, line_sp=1.15, sp_after=5)
card(s, 0.82+cw, y2, cw, 1.95, accent=C["ar"])
txt(s, 1.02+cw, y2+0.14, cw-0.4, 0.3, [[("Baseline — AR (snake ordering)",13,C["ar"],True,False,SANS)]])
txt(s, 1.02+cw, y2+0.5, cw-0.4, 1.35, [
    "• 비트를 snake 순서로 하나씩 예측: p(b)=∏ p(b_k | b_<k, y).",
    "• chain rule 을 정확히 따르므로 인수분해 오차가 없다.",
    "• 생성에 100 번의 순차 forward 가 필요하고, 병렬화되지 않는다."], size=12, line_sp=1.15, sp_after=5)
box(s, 0.62, y2+2.1, 12.1, 0.85, [
    "Flow model 은 100 개 픽셀을 매 스텝 병렬로 갱신한다. 관심은 적은 스텝으로 AR 의 품질에 도달하는지다. "
    "아래에서 두 접근을 본다 — CFM 연속 완화, DFM 이산 상태 유지."])

# ======================================================================
# 2. CFM framework
# ======================================================================
s = slide()
y = kicker_title(s, "Approach 1 · Framework", "CFM — 기존 vector field 프레임워크를 그대로 가져온다",
    "이진 마스크를 연속값으로 완화한 뒤, 표준 Continuous Flow Matching 을 수정 없이 적용한다.", num="2 / 10")
fw=2.9
for i,(a,b,c) in enumerate([("SOURCE","x0 ~ N(0, I)","표준 가우시안 노이즈"),
    ("TARGET","x1 ~ q(x1)","GT 데이터셋. {0,1}→{−1,+1}"),
    ("PATH","직선 보간","두 점을 직선으로 잇는다"),
    ("LEARN","velocity v","그 직선의 속도를 회귀")]):
    x=0.62+i*(fw+0.1); card(s,x,y,fw,0.9)
    txt(s,x+0.16,y+0.11,fw-0.3,0.25,[[(a,9.5,C["muted"],False,False,MONO)]])
    txt(s,x+0.16,y+0.32,fw-0.3,0.3,[[(b,12.5,C["ink"],False,False,MONO)]])
    txt(s,x+0.16,y+0.58,fw-0.3,0.3,[[(c,11,C["ink2"],False,False,SANS)]],line_sp=1.05)
y2=y+1.1; cw=6.0
card(s,0.62,y2,cw,1.85,accent=C["cfm"])
txt(s,0.82,y2+0.14,cw-0.4,0.3,[[("경로와 학습 대상",12.5,C["ink"],True,False,SANS)]])
txt(s,0.82,y2+0.5,cw-0.4,0.5,[[("x_t = (1−t)·x0 + t·x1        u_t = x1 − x0",13,C["ink"],False,False,MONO)]])
txt(s,0.82,y2+0.95,cw-0.4,0.8,[[("모델은 v(x_t,t,y) 로 이 속도를 예측하고 MSE 회귀로 학습한다. 이진 구조에 대한 별도 처리는 없다 — 이미지 생성 CFM 과 동일.",11.5,C["ink2"],False,False,SANS)]],line_sp=1.15)
card(s,0.82+cw,y2,cw,1.85,accent=C["cfm"])
txt(s,1.02+cw,y2+0.14,cw-0.4,0.3,[[("Inference — 적분한 뒤 이산화",12.5,C["ink"],True,False,SANS)]])
txt(s,1.02+cw,y2+0.5,cw-0.4,1.2,[
    "• x0 ~ N(0,I) 에서 Euler 적분으로 t: 0→1.",
    "• 도착한 x1 은 연속값 — 아직 이진 마스크가 아니다.",
    "• 마지막에 0 을 기준으로 thresholding: bits = (x > 0)."],size=11.5,line_sp=1.15,sp_after=4)
box(s,0.62,y2+2.0,12.1,1.15,[
    [("이 CFM 은 이진 대상을 연속 완화로 다루는 접근의 baseline 이다. 핵심 문제는 이산성 우회에 있다: 데이터는 {0,1}^100 격자 위에 있는데 "
     "연속공간에서 풀고 마지막에 thresholding 으로 되돌린다. 그 과정에서 |x|(per-bit confidence)는 버려진다.",12.5,C["ink2"],False,False,SANS)],
    [("참고: 이 구현이 결정론적인 것은 sigma_min=0 이라 경로에 노이즈가 없기 때문 — CFM 자체 성질이 아니라 설계 선택. one-to-many/다양성은 이 baseline 의 축이 아니다.",11,C["muted"],False,False,SANS)]],
    accent="dfm")

# ======================================================================
# 3. CFM code
# ======================================================================
s = slide()
y = kicker_title(s, "Approach 1 · Code", "CFM — Training & Sampling",
    "flow_matching.py — cfm_loss / sample_fm_cfg", num="3 / 10")
cw=6.0
codebox(s,0.62,y,cw,3.3,[
    "# x1: (B,100) in {0,1} → {−1,+1}",
    "x1 = xb * 2 - 1",
    "",
    "# 1) 노이즈와 시간",
    "x0 = torch.randn_like(x1)   # N(0,I)",
    "t  = torch.rand(B)          # U(0,1)",
    "",
    "# 2) 직선 경로 위의 점",
    "x_t = (1 - t) * x0 + t * x1",
    "",
    "# 3) 회귀 타겟 = 직선의 속도",
    "u_t = x1 - x0",
    "",
    "# 4) MSE 회귀",
    "v = model(x_t, t, y)        # out: R^100",
    "loss = ((v - u_t) ** 2).mean()"], title="Training")
codebox(s,0.82+cw,y,cw,3.3,[
    "# Euler ODE 적분, t: 0 → 1",
    "x  = torch.randn(B, 100)    # x0 ~ N(0,I)",
    "dt = 1.0 / N",
    "",
    "for i in range(N):",
    "    t = i * dt",
    "    v = model(x, t, y)      # velocity",
    "    x = x + v * dt          # clamp 없음",
    "",
    "# 연속값 → 이진 마스크",
    "bits = (x > 0).long()       # thresholding"], title="Sampling")
box(s,0.62,y+3.5,12.1,0.85,[[("CFM 입력/출력:  x_t ∈ R^100 (연속), t, y  →  v ∈ R^100 velocity (확률 아님).  Loss = MSE.  Inference 끝에 thresholding 필요.  NFE = N (병렬).",12.5,C["ink2"],False,False,SANS)]],accent="dfm")

# ======================================================================
# 4. DFM framework (+ g=1 diagram)
# ======================================================================
s = slide()
y = kicker_title(s, "Approach 2 · Framework", "DFM — 이산 상태를 유지한 채, 확률분포의 velocity field 를 학습",
    "각 픽셀의 vocab 을 {0,1} 로 정의하고, 모델이 상태의 확률분포를 직접 예측하게 한다.", num="4 / 10")
fw=2.9
for i,(a,b,c) in enumerate([("VOCAB","V = {0, 1}","픽셀당 2 심볼. 상태는 끝까지 이산"),
    ("SOURCE","p0 = Uniform","균등분포에서 출발"),
    ("TARGET","p1 = q(x1)","GT 데이터셋"),
    ("LEARN","확률의 velocity","p0 → p1 로 나르는 rate")]):
    x=0.62+i*(fw+0.1); card(s,x,y,fw,0.86)
    txt(s,x+0.16,y+0.1,fw-0.3,0.25,[[(a,9.5,C["muted"],False,False,MONO)]])
    txt(s,x+0.16,y+0.31,fw-0.3,0.3,[[(b,12.5,C["ink"],False,False,MONO)]])
    txt(s,x+0.16,y+0.57,fw-0.3,0.3,[[(c,10.5,C["ink2"],False,False,SANS)]],line_sp=1.02)
# g=1 diagram card
dy=y+1.02
card(s,0.62,dy,12.1,1.5,accent=C["dfm"])
txt(s,0.82,dy+0.12,11.6,0.3,[[("모델이 예측하는 것 (V=2) — 픽셀마다 확률을 내고, 거기서 {0,1} 을 뽑는다",12.5,C["ink"],True,False,SANS)]])
# prob 2x2
gx, gy, cell = 3.4, dy+0.5, 0.44
txt(s, gx-0.05, gy-0.24, 2*cell+0.1, 0.2, [[("픽셀별 P(=금속)",10,C["muted"],False,False,SANS)]], align=PP_ALIGN.CENTER)
for (cx,cy),p in zip([(0,0),(1,0),(0,1),(1,1)], ["0.90","0.20","0.80","0.70"]):
    rect(s, gx+cx*cell, gy+cy*cell, cell, cell, fill=C["sub"], line=C["line"], lw=0.75)
    txt(s, gx+cx*cell, gy+cy*cell, cell, cell, [[(p,12,C["ink"],False,False,MONO)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
arrow(s, gx+2*cell+0.15, gy+cell-0.2)
txt(s, gx+2*cell+0.05, gy-0.24, 0.8, 0.2, [[("샘플",10,C["muted"],False,False,SANS)]], align=PP_ALIGN.CENTER)
rx = gx+2*cell+0.85
patch(s, rx, gy, cell, ['m','e','m','m'], outline=C["dfm"])
txt(s, rx-0.1, gy+2*cell+0.02, 2*cell+0.2, 0.2, [[("뽑힌 마스크",10,C["muted"],False,False,SANS)]], align=PP_ALIGN.CENTER)
txt(s, 6.2, gy-0.05, 3.6, 1.0, [
    [("이 픽셀별 확률 p₁|t(x₁|x_t) 이",11.5,C["ink2"],False,False,SANS)],
    [("아래 rate 를 결정한다.",11.5,C["ink2"],False,False,SANS)],
    [("(V=16 으로 묶는 법은 6 장)",10.5,C["muted"],False,False,SANS)]], line_sp=1.15, sp_after=3)
# equations row
ey=dy+1.65; ew=3.9
eqs=[("1 · 조건부 경로","p_t(x|x1) = t·δ(x,x1) + (1−t)/V","확률 t 로 정답 유지, 아니면 uniform 재추출"),
     ("2 · Marginalization","p_t(x) = E_x1~q [ p_t(x|x1) ]","균등분포 → 데이터분포 확률경로"),
     ("3 · Rate (denoiser 로 복원)","u_t = λ_t·[ p1|t(v|z) − δ(v,z) ]","모델은 사후확률만 예측, rate 는 거기서")]
for i,(h,eq,d) in enumerate(eqs):
    x=0.62+i*(ew+0.1); card(s,x,ey,ew,1.5,accent=C["dfm"])
    txt(s,x+0.16,ey+0.12,ew-0.3,0.3,[[(h,11.5,C["ink"],True,False,SANS)]])
    rect(s,x+0.16,ey+0.5,ew-0.32,0.42,fill=C["bg"],line=C["line2"],lw=0.5)
    txt(s,x+0.16,ey+0.5,ew-0.32,0.42,[[(eq,10.5,C["ink"],False,False,MONO)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+0.16,ey+1.0,ew-0.3,0.45,[[(d,10.5,C["ink2"],False,False,SANS)]],line_sp=1.1)

# ======================================================================
# 5. DFM code
# ======================================================================
s = slide()
y = kicker_title(s, "Approach 2 · Code", "DFM — Training & Sampling",
    "flow_matching.py — dfm_loss / sample_dfm_cfg", num="5 / 10")
cw=6.0
codebox(s,0.62,y,cw,3.15,[
    "t = torch.rand(B)               # U(0,1)",
    "",
    "# 1) 조건부 경로에서 x_t 샘플",
    "keep = torch.rand(B, D) < t[:,None]",
    "rnd  = torch.randint(0, 2, (B, D)).float()",
    "x_t  = torch.where(keep, x1, rnd)",
    "",
    "# 2) 깨끗한 x1 의 사후확률 예측",
    "logits = model(x_t*2-1, t, y)   # R^100",
    "",
    "# 3) BCE — 타겟은 원본 x1",
    "loss = F.binary_cross_entropy_with_logits(",
    "           logits, x1)"], title="Training")
codebox(s,0.82+cw,y,cw,3.15,[
    "x = torch.randint(0, 2, (B, D)).float()",
    "for i in range(N):",
    "    t  = i / N",
    "    p1 = sigmoid(model(x*2-1, t, y))",
    "    r  = (p1>.5) if last else (rand<p1)",
    "    a  = 1 / (N - i)     # = (t'−t)/(1−t)",
    "    x  = torch.where(rand<a, r, x)",
    "return x.long()          # 이미 이진"], title="Sampling — CTMC 유한전이")
# compare table
txt(s,0.62,y+3.35,12,0.3,[[("입력 / 출력 — CFM 과의 대비",12.5,C["ink"],True,False,SANS)]])
table(s,0.62,y+3.7,12.1,["","입력","출력","Loss","Inference 끝"],
    [["CFM","x_t ∈ R^100 (연속)","v ∈ R^100 velocity","MSE","thresholding 필요"],
     ["DFM","x_t ∈ {0,1}^100 (이산)","sigmoid = p1|t(x1|x_t,y)","BCE","이미 이진"]],
    colw=[1.1,2.6,2.8,1.0,2.0], fs=11, row_h=0.4)

# ======================================================================
# 6. Patchify diagram
# ======================================================================
s = slide()
y = kicker_title(s, "Patchify · 그림으로", "Patchify 를 그림으로",
    "위: 픽셀을 묶어 토큰을 만드는 방법. 아래: 그로 인해 달라지는 점. 다음 장에서 코드와 함께 정리.", num="6 / 10")
# grouping strip
card(s,0.62,y,12.1,0.92,accent=C["dfm"])
txt(s,0.82,y+0.1,11.6,0.25,[[("2×2 픽셀을 토큰 1 개로 묶는다  —  2×2 흑백 배치는 2⁴ = 16 가지 → 각 배치를 토큰 하나로 (vocab 16)",12,C["ink"],True,False,SANS)]])
gcell=0.36; gx=4.6; gyy=y+0.44
patch(s,gx,gyy,gcell,['m','e','m','m'],outline=C["dfm"])
arrow(s,gx+2*gcell+0.15,gyy+gcell-0.2)
rect(s,gx+2*gcell+0.75,gyy+0.05,1.2,0.5,fill=C["dfmbg"],line=C["dfm"],lw=1)
txt(s,gx+2*gcell+0.75,gyy+0.05,1.2,0.5,[[("토큰",13,C["dfm"],True,False,SANS)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
txt(s,gx+2*gcell+0.55,gyy+0.58,1.6,0.2,[[("16 가지 중 하나",10,C["muted"],False,False,SANS)]],align=PP_ALIGN.CENTER)
# g=1 / g=4
ry=y+1.05; cw=6.0
card(s,0.62,ry,cw,1.85)
txt(s,0.82,ry+0.1,cw-0.4,0.25,[[("g=1 · 픽셀마다 확률  (4 장에서 본 방식)",12,C["ink"],True,False,SANS)]])
pcell=0.4; px=1.1; py=ry+0.5
for (cx,cy),p in zip([(0,0),(1,0),(0,1),(1,1)],["0.90","0.20","0.80","0.70"]):
    rect(s,px+cx*pcell,py+cy*pcell,pcell,pcell,fill=C["sub"],line=C["line"],lw=0.75)
    txt(s,px+cx*pcell,py+cy*pcell,pcell,pcell,[[(p,11,C["ink"],False,False,MONO)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
arrow(s,px+2*pcell+0.1,py+pcell-0.2)
patch(s,px+2*pcell+0.7,py,pcell,['m','e','m','m'],outline=C["dfm"])
txt(s,0.82,ry+1.5,cw-0.4,0.3,[[("픽셀마다 확률을 정하고, 각 픽셀을 금속/빈칸으로 따로 뽑는다",10.5,C["ink2"],False,False,SANS)]])
card(s,0.82+cw,ry,cw,1.85,accent=C["dfm"])
txt(s,1.02+cw,ry+0.1,cw-0.4,0.25,[[("g=4 · 배치마다 확률",12,C["ink"],True,False,SANS)]])
# bars
bx=1.02+cw+0.15; baseY=ry+1.35; heights=[.02,.02,.02,.02,.02,.15,.02,.02,.02,.28,.02,.11,.02,.9,.02,.47]
for i,h in enumerate(heights):
    col = C["dfm"] if i==13 else C["muted"]
    r=rect(s,bx+i*0.14,baseY-h,0.1,h,fill=col)
    if i!=13: r.fill.fore_color.rgb=clr(C["muted"])
txt(s,bx+13*0.14-0.15,baseY-0.9-0.22,0.4,0.2,[[("0.42",9.5,C["dfm"],False,False,MONO)]],align=PP_ALIGN.CENTER)
arrow(s,bx+16*0.14+0.05,ry+0.55,0.4,16)
patch(s,bx+16*0.14+0.5,ry+0.4,0.28,['m','e','m','m'])
txt(s,1.02+cw+0.2,ry+1.5,cw-0.4,0.3,[[("16 가지 배치 전체에 확률을 정하고, 그중 하나를 통째로 뽑는다",10.5,C["ink2"],False,False,SANS)]])
# note
txt(s,0.62,ry+1.95,12.1,0.5,[[("두 방법 모두 같은 배치가 나올 수 있다. 차이는 확률을 어디에 두는가다 — g=1 은 픽셀을 따로 뽑아 픽셀 사이 관계를 표현 못 하고, g=4 는 16 가지 배치에 직접 확률을 줘 그 관계를 담는다.",11,C["ink2"],False,False,SANS)]],align=PP_ALIGN.CENTER,line_sp=1.1)
# collision / noising
by=ry+2.5; cw2=6.0
card(s,0.62,by,cw2,1.35)
txt(s,0.82,by+0.1,cw2-0.4,0.25,[[("collision — 노이즈가 원래 값과 우연히 같을 확률",12,C["ink"],True,False,SANS)]])
oval(s,1.3,by+0.5,0.7,0.7,fill=C["sub"],line=C["line"],lw=1)
txt(s,1.3,by+0.5,0.7,0.7,[[("1/2",12,C["dfm"],True,False,MONO)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
txt(s,1.25,by+0.4,0.8,0.15,[[("V=2",9.5,C["muted"],False,False,SANS)]],align=PP_ALIGN.CENTER)
oval(s,3.0,by+0.5,0.7,0.7,fill=C["sub"],line=C["line"],lw=1)
txt(s,3.0,by+0.5,0.7,0.7,[[("1/16",11,C["dfm"],True,False,MONO)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
txt(s,2.95,by+0.4,0.8,0.15,[[("V=16",9.5,C["muted"],False,False,SANS)]],align=PP_ALIGN.CENTER)
txt(s,4.2,by+0.55,2.3,0.7,[[("묶을수록 낮아진다 → 노이즈를 줬을 때 실제로 더 많이 바뀐다",10.5,C["ink2"],False,False,SANS)]],line_sp=1.1,anchor=MSO_ANCHOR.MIDDLE)
card(s,0.82+cw2,by,cw2,1.35)
txt(s,1.02+cw2,by+0.1,cw2-0.4,0.25,[[("노이징 단위 — 비트마다 vs 패치 통째로",12,C["ink"],True,False,SANS)]])
nx=1.4+cw2; ncell=0.26
patch(s,nx,by+0.55,ncell,['m','x','x','m'])
txt(s,nx-0.15,by+1.08,2*ncell+0.3,0.2,[[("비트마다",9.5,C["muted"],False,False,SANS)]],align=PP_ALIGN.CENTER)
patch(s,nx+1.4,by+0.55,ncell*0.7,['m','m','m','m'])
txt(s,nx+1.9,by+0.6,0.3,0.3,[[("or",10,C["muted"],False,False,SANS)]])
patch(s,nx+2.3,by+0.55,ncell*0.7,['x','x','x','x'])
txt(s,nx+1.3,by+1.08,2.0,0.2,[[("패치 통째로",9.5,C["muted"],False,False,SANS)]],align=PP_ALIGN.CENTER)

# ======================================================================
# 7. Patchify text/code
# ======================================================================
s = slide()
y = kicker_title(s, "Approach 2+ · Vocab merging", "Patchify — 픽셀을 묶어 vocab 을 키운다",
    "DFM 을 V=2 로 그대로 쓰면 학습이 잘 되지 않는다. 원인은 collision 이다.", num="7 / 10")
cw=6.0
card(s,0.62,y,cw,2.5)
txt(s,0.82,y+0.12,cw-0.4,0.3,[[("문제 — collision",12.5,C["ink"],True,False,SANS)]])
txt(s,0.82,y+0.48,cw-0.4,0.6,[[("uniform 재추출이 원래 값과 같을 확률이 1/V. 실제 flip 확률 = (1−t)·(V−1)/V.",11.5,C["ink2"],False,False,SANS)]],line_sp=1.1)
table(s,0.82,y+1.05,cw-0.4,["V","collision","t=0 실제 flip"],
    [["2","0.500","0.500"],["16","0.063","0.938"],["1024","0.001","0.999"]],
    colw=[1,1,1.4],fs=11,row_h=0.32,hi={(0,2)})
txt(s,0.82,y+2.1,cw-0.4,0.35,[[("V=2 는 t=0 에서도 절반만 뒤집혀 gradient 가 잘 서지 않는다.",11,C["ink2"],False,False,SANS)]],line_sp=1.05)
codebox(s,0.82+cw,y,cw,2.5,[
    "# 2×2 공간 패치로 재배열 후 묶는다",
    "fwd, inv = make_spatial_patch_indices(H,W,2,2)",
    "X = X[:, fwd]",
    "tok = bits_to_tokens(x1, group_size=4)",
    "",
    "# 노이징: 토큰 단위 uniform 재추출",
    "tok_t = torch.where(keep, tok,",
    "          torch.randint(0,16,(B,T)))",
    "",
    "# Loss: BCE → Cross-entropy",
    "loss = F.cross_entropy(",
    "    model(tok_t,t,y).reshape(-1,16),",
    "    tok.reshape(-1))"], title="해법 — 2×2 패치를 1 토큰으로 (g=4, V=16)")
# 3 effects
ey=y+2.7; ew=3.9
effs=[("1 · 출력 형태","g=1 은 비트별 베르누이 100 개(독립). g=4 는 토큰당 16-way categorical 로 묶인 4 비트의 결합분포를 표현."),
      ("2 · collision (측정됨)","1/2 → 1/16 로 낮아진다. t=0 실제 flip 0.500 → 0.938. 07-05 의 V=2 학습 부진과 직접 연결."),
      ("3 · 노이징 단위","g=1 은 비트마다 독립 → 패치 유지확률 t⁴. g=4 는 토큰 단위 → t. 유지/교체가 패치 단위로 일어난다.")]
for i,(h,d) in enumerate(effs):
    x=0.62+i*(ew+0.1); card(s,x,ey,ew,1.35,accent=C["dfm"])
    txt(s,x+0.16,ey+0.12,ew-0.3,0.3,[[(h,11.5,C["ink"],True,False,SANS)]])
    txt(s,x+0.16,ey+0.48,ew-0.3,0.8,[[(d,10.5,C["ink2"],False,False,SANS)]],line_sp=1.15)

# ======================================================================
# 8. Training curves
# ======================================================================
s = slide()
y = kicker_title(s, "Results · Training", "학습 곡선",
    "동일 백본 (Stable3DiT 76.8M · L=15 · d=512 · nhead=4 · dim_ff=768), 500 epochs · bs=128 · lr=1e-4.", num="8 / 10")
pw=6.0
for i,(cap) in enumerate(["train / val loss — CFM (MSE) · DFM V=2 (BCE) · DFM V=16 (CE)",
                          "val BitAcc / PatAcc — patchify (V=16) 의 효과가 드러나는 지점"]):
    x=0.62+i*(pw+0.2)
    r=rect(s,x,y,pw,3.2,fill=C["sub"],line=C["line"],lw=1,dash=True)
    txt(s,x,y+1.2,pw,0.3,[[("wandb capture",11,C["muted"],False,False,MONO)]],align=PP_ALIGN.CENTER)
    txt(s,x+0.5,y+1.55,pw-1.0,0.8,[[(cap,12,C["ink2"],False,False,SANS)]],align=PP_ALIGN.CENTER,line_sp=1.15)
txt(s,0.62,y+3.4,12.1,0.4,[[("CFG 용 체크포인트는 drop_prob 0.1 (조건을 10% 확률로 drop 해 uncond 도 함께 학습), 비교군은 0.0.",11,C["muted"],False,False,SANS)]])

# ======================================================================
# 9. Eval results
# ======================================================================
s = slide()
y = kicker_title(s, "Results · Evaluation", "평가 결과 — DFM + patchify 가 AR 을 앞선다",
    "주 지표는 notch-region MSE (논문 저자 코드). 공진 test 2,000 타겟, cfg=1.0, FM steps=10. "
    "notch MSE·valid@10 은 best-of-10, BitAcc·PatAcc 는 pass@1.", num="9 / 10")
table(s,0.62,y,12.1,["모델","NFE","BitAcc↑ (sanity)","PatAcc↑ (sanity)","notch MSE↓ (주 지표)","valid@10↑"],
    [["AR (snake)","100 순차","0.539","0.000","39.60","0.633"],
     ["CFM","10 병렬","0.730","0.005","60.64","0.398"],
     ["DFM V=2","10 병렬","0.809","0.141","44.98","0.517"],
     ["DFM V=16 (patchify)","10 병렬","0.831","0.115","36.53","0.591"]],
    colw=[2.4,1.4,1.7,1.7,2.0,1.4], fs=12, row_h=0.42,
    hi={(0,4),(0,1),(3,1)}, win={(3,2),(3,4),(0,5),(2,3)})
cy=y+2.0; cw=3.9
cards9=[("1 · DFM V=16 이 AR 을 앞선다","notch MSE 36.53 vs 39.60 — DFM 이 낮다. AR 은 100 회 순차, DFM 은 10 회 병렬. guidance 없이 얻은 결과."),
        ("2 · V=2 와 V=16 의 차이","V=2 는 44.98, V=16 은 36.53. V=2 는 AR(39.60)도 넘지 못한다. collision 의 영향이 성능차로."),
        ("3 · CFM (이산성 우회)","가장 높다 (60.64). 연속으로 풀고 thresholding 으로 되돌리는 손해. DFM 은 이산 유지.")]
for i,(h,d) in enumerate(cards9):
    x=0.62+i*(cw+0.1); ac=C["cfm"] if i==2 else C["dfm"]; card(s,x,cy,cw,1.3,accent=ac)
    txt(s,x+0.16,cy+0.12,cw-0.3,0.3,[[(h,11.5,C["ink"],True,False,SANS)]])
    txt(s,x+0.16,cy+0.48,cw-0.3,0.75,[[(d,10.5,C["ink2"],False,False,SANS)]],line_sp=1.12)
box(s,0.62,cy+1.45,12.1,1.15,[
    [("지표 참고. notch MSE 는 공진 구간에서만 잰 MSE (test 타겟의 74% 는 공진 없음). 지표 정의·코드는 논문을 따랐으나 측정 대상과 후보 예산(best-of-10)이 논문과 달라 논문 절대값(56.01)과 직접 비교 아님 — 목적은 AR vs FM 상대 비교.",11.5,C["ink2"],False,False,SANS)],
    [("valid@10 은 AR 이 약간 높다(0.633 vs 0.591) — 다양성은 AR 이 넓지만 best 후보 품질(notch MSE)은 DFM 이 앞선다. BitAcc·PatAcc 는 GT 일치율이라 one-to-many 에서 대안해를 벌주는 sanity 지표. AR 은 실제 100 스텝 생성으로 평가.",10.5,C["muted"],False,False,SANS)]],
    accent="dfm")

# ======================================================================
# 10. Steps
# ======================================================================
s = slide()
y = kicker_title(s, "Results · Steps", "스텝 수에 따른 후보 다양성과 notch MSE",
    "DFM V=16, guidance 없음 (cfg = 1.0), 공진 test 2,000 타겟, best-of-10.", num="10 / 10")
table(s,0.62,y,12.1,["NFE (steps)","평균 Hamming 거리","서로 다른 후보 / 10","BitAcc↑","PatAcc↑","notch MSE↓ (주 지표)","valid@10↑"],
    [["1","1.5","3.5","0.871","0.324","58.68","0.429"],
     ["2","10.2","8.0","0.850","0.182","36.06","0.590"],
     ["4","13.3","8.7","0.837","0.134","35.62","0.593"],
     ["10","14.6","9.0","0.831","0.115","36.53","0.591"],
     ["50","14.8","9.1","0.831","0.108","36.42","0.604"]],
    colw=[1.4,1.9,1.8,1.2,1.2,2.0,1.3], fs=12, row_h=0.44, hi={(0,0),(0,5)})
txt(s,0.62,y+2.9,12.1,0.6,[[("Hamming 거리 = 같은 타겟에 대해 생성한 후보 10 개가 서로 평균 몇 비트 다른가 (100 비트 중). 9 장의 결과는 steps = 10 에서 잰 값이다.",11.5,C["muted"],False,False,SANS)]],line_sp=1.15)

prs.save("/hai/home/lsh/antenna/year_hai/src/Flow_matching/docs/fm-meeting-deck.pptx")
print("total slides:", len(prs.slides._sldIdLst))
